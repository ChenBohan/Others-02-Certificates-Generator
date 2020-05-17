# -*- coding:utf-8 -*-

from pptx import Presentation
import os
import sys
from tkinter import filedialog, messagebox, Tk, Label, Button, StringVar, mainloop
from datetime import datetime
import logging
import tkinter as tk # Python 3.x
import tkinter.scrolledtext as ScrolledText
import pinyin
import argparse

txt_chn_path = ''
ppt_path = ''
ppt_filename = 'certificates'

class TextHandler(logging.Handler):
# This class allows you to log to a Tkinter Text or ScrolledText widget
# Adapted from Moshe Kaplan: https://gist.github.com/moshekaplan/c425f861de7bbf28ef06

    def __init__(self, text):
        # run the regular Handler __init__
        logging.Handler.__init__(self)
        # Store a reference to the Text it will log to
        self.text = text

    def emit(self, record):
        msg = self.format(record)
        def append():
            self.text.configure(state='normal')
            self.text.insert(tk.END, msg + '\n')
            self.text.configure(state='disabled')
            # Autoscroll to the bottom
            self.text.yview(tk.END)
        # This is necessary because we can't modify the Text from other threads
        self.text.after(0, append)


sys.stdout = open("./log.txt", "w")
print ("test sys.stdout")

def browse_txt_button():
    global show_txt_path
    global txt_chn_path
    txt_chn_path = filedialog.askopenfilename(initialdir = ".",title = "Select file",filetypes = (("ppt files","*.txt"),("all files","*.*")))
    if txt_chn_path:
        logging.info(now.strftime("%H:%M:%S") + " 成功导入名单：" + txt_chn_path)
        show_txt_path.set(os.path.basename(txt_chn_path))
        print(range(len(txt_chn_path)))
    else:
        logging.info(now.strftime("%H:%M:%S") + " 未导入名单")
        show_txt_path.set("未导入名单，请重试")
        
def browse_ppt_button():
    global show_ppt_path
    global ppt_path
    ppt_path = filedialog.askopenfilename(initialdir = ".",title = "Select file",filetypes = (("txt files","*.pptx"),("all files","*.*")))
    if ppt_path:
        logging.info(now.strftime("%H:%M:%S") + " 成功导入PPT：" + ppt_path)
        show_ppt_path.set(os.path.basename(ppt_path))
    else:
        show_ppt_path.set("未导入PPT， 请重试")

def browse_rename_button():
    global show_rename_path
    global txt_chn_path
    global ppt_filename
    rename_path = filedialog.askdirectory()
    show_txt_path.set(rename_path)
    print(rename_path)

    with open(txt_chn_path, 'r', encoding='UTF-8') as f:
        name_list = f.readlines()

    for i in range(len(txt_chn_path) - 1):
        if i <= 8:
            rename_file_name = "{}/{}_0{}.png".format(rename_path, ppt_filename, i+1)
            if os.path.isfile(rename_file_name):
                os.rename(rename_file_name, "{}/{}.png".format(rename_path, name_list[i].strip()))
                logging.info(now.strftime("%H:%M:%S") + " 成功将 {} 重命名为 {}".format("{}/{}_0{}.png".format(rename_path, ppt_filename, i+1), "{}.png".format(name_list[i].strip())))
            else:
                continue
        else:
            rename_file_name = "{}/{}_{}.png".format(rename_path, ppt_filename, i+1)
            if os.path.isfile(rename_file_name):
                os.rename(rename_file_name, "{}/{}.png".format(rename_path, name_list[i].strip()))
                logging.info(now.strftime("%H:%M:%S") + " 成功将 {} 重命名为 {}".format("{}/{}_0{}.png".format(rename_path, ppt_filename, i+1), "{}.png".format(name_list[i].strip())))
            else:
                continue
    messagebox.showinfo("提示", "重命名完成")

def replace_name_ppt(name_list):
    global ppt_path
    global ppt_filename
    global generate_log
    
    prs = Presentation(ppt_path)
    i = 0
    if len(prs.slides) >= len(name_list):
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == 'xxx':
                            if i >= len(name_list):
                                prs.save('{}.pptx'.format(ppt_filename))
                                messagebox.showinfo("提示", "生成成功（保存在当前目录下的'certificates.pptx'）")
                                return
                            current_name = name_list[i]
                            run.text = current_name
                            i = i + 1
                            logging.info(now.strftime("%H:%M:%S") + " 成功生成第{}张PPT：{}".format(i, current_name))
        prs.save('{}.pptx'.format(ppt_filename))
        messagebox.showinfo("提示", "生成成功（保存在当前目录下的'certificates.pptx'）")
    else:
        logging.info(now.strftime("%H:%M:%S") + " PPT页数不足：仅{}页，需{}页".format(len(prs.slides), len(name_list)))
    
def export_chn_button():
    global txt_chn_path
    with open(txt_chn_path, 'r', encoding='UTF-8') as f:
        name_list = f.readlines()
    name_list = [name.strip('\n') for name in name_list]
    replace_name_ppt(name_list)
    
def export_eng_button():
    global txt_chn_path
    with open(txt_chn_path, 'r', encoding='UTF-8') as f:
        name_list_chn = f.readlines()
    name_list_chn = [name.strip('\n') for name in name_list_chn]
    name_list_eng = []
    for chn_name in name_list_chn:
        eng_name = toPinyin(chn_name, True).strip()
        name_list_eng.append(eng_name)
        logging.info(now.strftime("%H:%M:%S") + " 成功将 {} 转换为 {}".format(chn_name, eng_name))
    replace_name_ppt(name_list_eng)

def isEnglish(s):
    try:
        s.encode('ascii')
    except UnicodeEncodeError:
        return False
    else:
        return True

def toPinyin(s, rearrange=True):
    s = s.strip()
    if isEnglish(s):
        return s
    chars = []
    if rearrange:
        firstname=s[1:]
        lastname=s[0]
        chars.append(pinyin.get(firstname, format="strip", delimiter="").strip("-"))
        chars.append(pinyin.get(lastname, format="strip", delimiter="-"))
    else:
        chars+= pinyin.get(s, format="strip", delimiter=" ").split()
    return " ".join(map(lambda x: x.capitalize(), chars))

row = 0
root = Tk()
root.after(1000)

show_txt_path = StringVar()
lbl1 = Label(master=root, text="1.导入名单（格式见 示例.txt）：")
lbl1.grid(row=row, column=0)
lbl_txt_path = Label(master=root, textvariable=show_txt_path)
lbl_txt_path.grid(row=row, column=1)
button_txt = Button(text="选择名单", command=browse_txt_button)
button_txt.grid(row=row, column=2)
row = row + 1

show_ppt_path = StringVar()
lbl3 = Label(master=root, text="2.导入pptx模板（名字预填小写无空格'xxx'，先复制好足够的页数）：")
lbl3.grid(row=row, column=0)
lbl4 = Label(master=root, textvariable=show_ppt_path)
lbl4.grid(row=row, column=1)
button3 = Button(text="选择PPT", command=browse_ppt_button)
button3.grid(row=row, column=2)
row = row + 1

generate_log = StringVar()
lbl5 = Label(master=root, text="3.生成PPT")
lbl5.grid(row=row, column=0)
lbl_gen_log = Label(master=root, textvariable=generate_log)
lbl_gen_log.grid(row=row, column=1)
button_gen_chn = Button(text="生成中文名奖状", command=export_chn_button)
button_gen_chn.grid(row=row, column=2)
button_gen_eng = Button(text="生成英文名奖状", command=export_eng_button)
button_gen_eng.grid(row=row, column=3)
row = row + 1

lbl12 = Label(master=root, text="4.手动将PPT逐页导出为图片（图片名应为'PPT名_序号'）")
lbl12.grid(row=row, column=0)
row = row + 1

show_rename_path = StringVar()
lbl6 = Label(master=root, text="5.图片名改为中文名")
lbl6.grid(row=row, column=0)
lbl7 = Label(master=root, textvariable=show_rename_path)
lbl7.grid(row=row, column=1)
button5 = Button(text="选择图片所在文件夹", command=browse_rename_button)
button5.grid(row=row, column=2)
row = row + 1

# Add text widget to display logging info
st = ScrolledText.ScrolledText(state='disabled')
st.configure(font='TkFixedFont')
st.grid(column=0, row=row, sticky='w', columnspan=4)

# datetime object containing current date and time
now = datetime.now()

# Create textLogger
text_handler = TextHandler(st)

# Logging configuration
logging.basicConfig(filename='test.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s')

# Add the handler to logger
logger = logging.getLogger()
logger.addHandler(text_handler)

logging.info("------运行日志------")
logging.info(now.strftime("%H:%M:%S") + " 奖状生成器初始化成功")

mainloop()



